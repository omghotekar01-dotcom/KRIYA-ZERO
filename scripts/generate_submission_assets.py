from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib.colors import HexColor
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.units import mm
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "submission"
OUT.mkdir(parents=True, exist_ok=True)

PPTX_PATH = OUT / "KRIYA_ZERO_FINAL_DECK.pptx"
PDF_PATH = OUT / "KRIYA_ZERO_FINAL_IDEA.pdf"

NAVY = RGBColor(11, 22, 40)
INK = RGBColor(22, 32, 48)
TEAL = RGBColor(0, 173, 181)
MINT = RGBColor(220, 249, 246)
WHITE = RGBColor(255, 255, 255)
SOFT = RGBColor(242, 246, 249)
GRAY = RGBColor(99, 115, 129)
RED = RGBColor(217, 70, 70)
GREEN = RGBColor(45, 156, 110)


def add_text(slide, text, x, y, w, h, size=24, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_round_rect(slide, x, y, w, h, fill=WHITE, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or fill
    if line is None:
        shp.line.fill.background()
    return shp


def add_badge(slide, text, x, y, w, fill=TEAL, color=WHITE):
    add_round_rect(slide, x, y, w, 0.38, fill=fill)
    add_text(slide, text, x + 0.08, y + 0.03, w - 0.16, 0.28, 10, color, True, align=PP_ALIGN.CENTER)


def add_header(slide, kicker, title, subtitle=None, dark=False):
    c1 = WHITE if dark else NAVY
    c2 = RGBColor(203, 216, 225) if dark else GRAY
    add_text(slide, kicker.upper(), 0.7, 0.45, 5.5, 0.35, 11, TEAL, True)
    add_text(slide, title, 0.7, 0.86, 12.0, 0.8, 30, c1, True)
    if subtitle:
        add_text(slide, subtitle, 0.7, 1.62, 11.6, 0.58, 15, c2)


def set_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_footer(slide, n, dark=False):
    color = RGBColor(169, 184, 197) if dark else RGBColor(138, 151, 164)
    add_text(slide, f"KRIYA ZERO  •  XYRO  •  iQOO Pune 2026", 0.7, 7.14, 6.2, 0.22, 9, color)
    add_text(slide, f"{n:02d}", 12.2, 7.14, 0.4, 0.22, 9, color, True, align=PP_ALIGN.RIGHT)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_badge(s, "SMART EDUCATION", 0.75, 0.72, 1.75)
    add_text(s, "KRIYA ZERO", 0.75, 1.45, 6.4, 0.9, 42, WHITE, True)
    add_text(s, "The One-Demonstration Physical Skill Compiler", 0.75, 2.28, 6.9, 0.78, 23, RGBColor(211, 223, 232), True)
    add_text(s, "Show it once. Learn it anywhere. Prove you can do it.", 0.75, 3.23, 6.5, 0.65, 18, RGBColor(181, 198, 210))
    add_round_rect(s, 8.2, 1.05, 4.3, 4.75, fill=RGBColor(18, 34, 56), line=RGBColor(42, 64, 86))
    add_text(s, "DEMONSTRATION", 8.6, 1.6, 3.4, 0.35, 12, TEAL, True, align=PP_ALIGN.CENTER)
    add_text(s, "↓", 9.95, 2.1, 0.7, 0.5, 26, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(s, "EXECUTABLE\nPROCEDURAL GRAPH", 8.65, 2.65, 3.25, 0.95, 18, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(s, "↓", 9.95, 3.72, 0.7, 0.5, 26, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(s, "LIVE VERIFIER", 8.65, 4.34, 3.25, 0.45, 18, TEAL, True, align=PP_ALIGN.CENTER)
    add_text(s, "Android-first • phone-native • local-first", 8.65, 5.14, 3.25, 0.38, 11, RGBColor(170, 190, 204), align=PP_ALIGN.CENTER)
    add_text(s, "Team XYRO  •  Om Ghotekar  •  VIT Pune", 0.75, 6.55, 6.5, 0.35, 12, RGBColor(170, 190, 204))
    add_footer(s, 1, dark=True)

    # Slide 2
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_header(s, "The problem", "Watching is not the same as proving skill", "Practical education still struggles to verify whether a learner can correctly reproduce a physical procedure.")
    cards = [
        ("VIDEO ≠ VERIFICATION", "Tutorials show the procedure, but cannot reliably detect a learner's wrong, skipped or out-of-order physical step."),
        ("EXPERT TIME DOESN'T SCALE", "Labs, ITIs and vocational programs cannot place an expert beside every learner for every repetition."),
        ("ASSESSMENT IS OFTEN MANUAL", "Practical competence is difficult to measure consistently, especially across large cohorts and distributed training."),
    ]
    for i, (t, b) in enumerate(cards):
        x = 0.75 + i * 4.15
        add_round_rect(s, x, 2.45, 3.75, 3.3, fill=SOFT)
        add_text(s, f"0{i+1}", x + 0.22, 2.7, 0.65, 0.45, 15, TEAL, True)
        add_text(s, t, x + 0.22, 3.25, 3.25, 0.75, 17, NAVY, True)
        add_text(s, b, x + 0.22, 4.1, 3.25, 1.25, 13, GRAY)
    add_footer(s, 2)

    # Slide 3
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_header(s, "The invention", "A fresh demonstration becomes a reusable verifier", "KRIYA turns a single structured physical demonstration into a Skill Capsule: steps, dependencies, visual checkpoints and assessment logic.")
    stages = [
        ("1", "TEACH", "Expert demonstrates once while the phone observes the workspace."),
        ("2", "COMPILE", "KRIYA builds ordered checkpoints and an executable procedure graph."),
        ("3", "VERIFY", "A learner performs the task; wrong or skipped states are rejected."),
        ("4", "ASSESS", "Completion, corrections and sequence errors become evidence-backed metrics."),
    ]
    for i, (n, t, b) in enumerate(stages):
        x = 0.65 + i * 3.14
        add_round_rect(s, x, 2.55, 2.72, 3.45, fill=SOFT)
        add_round_rect(s, x + 0.18, 2.82, 0.52, 0.52, fill=TEAL)
        add_text(s, n, x + 0.18, 2.88, 0.52, 0.32, 13, WHITE, True, align=PP_ALIGN.CENTER)
        add_text(s, t, x + 0.18, 3.55, 2.3, 0.45, 17, NAVY, True)
        add_text(s, b, x + 0.18, 4.14, 2.3, 1.25, 12, GRAY)
        if i < 3:
            add_text(s, "→", x + 2.72, 3.82, 0.42, 0.45, 22, TEAL, True, align=PP_ALIGN.CENTER)
    add_footer(s, 3)

    # Slide 4
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_header(s, "Why KRIYA is different", "Not another video-to-SOP tool", "The defensible wedge is zero-authoring physical verification: the verifier is created from the demonstration itself.", dark=True)
    rows = [
        ("Video tutorials", "Show instructions", "No live execution verification"),
        ("Digital SOP / AR work instructions", "Guide predefined procedures", "Procedure/checkpoints are authored beforehand"),
        ("KRIYA ZERO", "Learns fresh checkpoints from demonstration", "Immediately becomes instructor + verifier"),
    ]
    y = 2.55
    for i, (a, b, c) in enumerate(rows):
        fill = RGBColor(24, 45, 69) if i < 2 else RGBColor(8, 91, 98)
        add_round_rect(s, 0.8, y, 11.75, 1.05, fill=fill, line=RGBColor(43, 65, 87))
        add_text(s, a, 1.05, y + 0.25, 3.0, 0.45, 15, WHITE, i == 2)
        add_text(s, b, 4.15, y + 0.22, 3.6, 0.55, 13, RGBColor(218, 228, 235))
        add_text(s, c, 7.9, y + 0.22, 4.2, 0.55, 13, RGBColor(218, 228, 235))
        y += 1.28
    add_text(s, "Safe novelty claim", 0.82, 6.45, 1.4, 0.28, 10, TEAL, True)
    add_text(s, "Learns previously unseen, visually observable structured procedures from one demonstration and creates reusable verification checkpoints without manual procedure authoring.", 2.25, 6.39, 10.0, 0.5, 11, RGBColor(190, 206, 218))
    add_footer(s, 4, dark=True)

    # Slide 5
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_header(s, "Technical approach", "AI/perception provides evidence; deterministic logic owns PASS/FAIL", "This split keeps the demo robust and makes future local models replaceable without rewriting the procedure engine.")
    pipeline = [
        ("CameraX", "Live rear-camera frames"),
        ("Visual evidence", "Learned checkpoint fingerprints + stabilization"),
        ("Skill Capsule", "States • dependencies • verifier evidence"),
        ("Verification Engine", "PASS • FAIL • sequence error • correction"),
        ("Assessment", "Completion • first-attempt accuracy • assistance"),
    ]
    for i, (t, b) in enumerate(pipeline):
        x = 0.55 + i * 2.5
        add_round_rect(s, x, 2.55, 2.12, 2.55, fill=MINT if i in (1, 2, 3) else SOFT)
        add_text(s, t, x + 0.16, 2.9, 1.8, 0.55, 15, NAVY, True, align=PP_ALIGN.CENTER)
        add_text(s, b, x + 0.16, 3.62, 1.8, 0.88, 11, GRAY, align=PP_ALIGN.CENTER)
        if i < len(pipeline) - 1:
            add_text(s, "→", x + 2.13, 3.48, 0.36, 0.45, 20, TEAL, True, align=PP_ALIGN.CENTER)
    add_round_rect(s, 1.25, 5.55, 10.7, 0.85, fill=SOFT)
    add_text(s, "Current working prototype", 1.52, 5.8, 2.2, 0.3, 12, TEAL, True)
    add_text(s, "Android • CameraX • learned visual checkpoints • five-frame stabilization • temporal skip detection • assessment • tests • CI/APK build", 3.6, 5.72, 7.95, 0.45, 12, INK)
    add_footer(s, 5)

    # Slide 6
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_header(s, "The 20-second jury moment", "Teach something fresh → make a mistake → KRIYA catches it", "The demo proves the core idea without relying on a pre-authored procedure.")
    add_round_rect(s, 0.8, 2.35, 3.6, 3.85, fill=SOFT)
    add_text(s, "TEACH", 1.08, 2.68, 1.3, 0.35, 13, TEAL, True)
    add_text(s, "Expert creates 3–5 visually distinct tabletop states in front of the phone.", 1.08, 3.25, 3.0, 1.0, 16, NAVY, True)
    add_text(s, "Checkpoint 01 learned\nCheckpoint 02 learned\nCheckpoint 03 learned", 1.08, 4.58, 2.8, 1.0, 12, GRAY)

    add_round_rect(s, 4.85, 2.35, 3.6, 3.85, fill=RGBColor(255, 244, 244))
    add_text(s, "INTENTIONAL ERROR", 5.13, 2.68, 2.2, 0.35, 13, RED, True)
    add_text(s, "Learner skips or performs the wrong state.", 5.13, 3.25, 2.95, 0.8, 16, NAVY, True)
    add_text(s, "EXECUTION DIVERGENCE\nExpected: Step 2\nObserved: wrong / later state", 5.13, 4.45, 2.9, 1.25, 12, RED, True)

    add_round_rect(s, 8.9, 2.35, 3.6, 3.85, fill=RGBColor(239, 250, 245))
    add_text(s, "CORRECT + VERIFY", 9.18, 2.68, 2.1, 0.35, 13, GREEN, True)
    add_text(s, "Learner fixes the state without restarting.", 9.18, 3.25, 2.95, 0.8, 16, NAVY, True)
    add_text(s, "✓ Correction verified\n✓ Procedure complete\nSkill report generated", 9.18, 4.45, 2.9, 1.25, 12, GREEN, True)
    add_footer(s, 6)

    # Slide 7
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_header(s, "Impact & scalability", "A practical-skill layer for education and workforce training", "Start with structured tabletop tasks; scale the same Skill Capsule architecture as local perception gets stronger.")
    domains = [
        ("Engineering labs", "Circuit assembly, instrumentation, practical procedures"),
        ("ITIs & vocational training", "Repeatable hands-on skill guidance and assessment"),
        ("Industrial workforce", "Onboarding, maintenance sequences, standardized practice"),
        ("Field technicians", "Guided procedures where an expert cannot be physically present"),
    ]
    for i, (t, b) in enumerate(domains):
        x = 0.75 + (i % 2) * 6.15
        y = 2.45 + (i // 2) * 1.75
        add_round_rect(s, x, y, 5.55, 1.42, fill=SOFT)
        add_text(s, t, x + 0.22, y + 0.23, 2.35, 0.42, 16, NAVY, True)
        add_text(s, b, x + 2.55, y + 0.2, 2.7, 0.72, 12, GRAY)
    add_round_rect(s, 0.75, 6.05, 11.7, 0.62, fill=MINT)
    add_text(s, "Feasible MVP scope: visually observable, structured tabletop procedures — deliberately narrow enough to work, broad enough to matter.", 1.0, 6.21, 11.2, 0.3, 12, INK, True, align=PP_ALIGN.CENTER)
    add_footer(s, 7)

    # Slide 8
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_header(s, "Why iQOO", "The phone is the product, not a remote control", "Camera, on-device perception and edge execution are central to KRIYA's value — ideal for an iQOO phone-first hackathon.", dark=True)
    add_badge(s, "ANDROID-FIRST", 0.85, 2.55, 1.55)
    add_badge(s, "PHONE CAMERA", 2.65, 2.55, 1.55, fill=RGBColor(33, 104, 126))
    add_badge(s, "LOCAL-FIRST", 4.45, 2.55, 1.45, fill=RGBColor(48, 86, 139))
    add_badge(s, "OFFLINE-CRITICAL PATH", 6.15, 2.55, 2.1, fill=RGBColor(88, 73, 145))
    add_text(s, "What exists today", 0.85, 3.42, 2.1, 0.4, 14, TEAL, True)
    add_text(s, "Working Android foundation, live CameraX input, learned checkpoints, deterministic verifier, sequence-error detection, assessment logic, persistence backend, unit tests and CI-generated APK.", 0.85, 3.9, 5.65, 1.25, 15, WHITE)
    add_text(s, "Next at the hackathon", 7.0, 3.42, 2.4, 0.4, 14, TEAL, True)
    add_text(s, "Strengthen viewpoint robustness with local image embeddings, add local narration, calibrate on the iQOO device, and polish the judge-facing demo — without weakening the reliable deterministic core.", 7.0, 3.9, 5.25, 1.25, 15, WHITE)
    add_text(s, "KRIYA ZERO", 0.85, 5.75, 4.0, 0.55, 28, WHITE, True)
    add_text(s, "Show it once. Learn it anywhere. Prove you can do it.", 0.85, 6.3, 8.2, 0.45, 17, RGBColor(198, 214, 225), True)
    add_text(s, "github.com/omghotekar01-dotcom/KRIYA-ZERO", 8.0, 6.34, 4.3, 0.35, 11, TEAL, True, align=PP_ALIGN.RIGHT)
    add_footer(s, 8, dark=True)

    prs.save(PPTX_PATH)


def build_pdf():
    styles = getSampleStyleSheet()
    title = ParagraphStyle("TitleK", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=24, leading=28, textColor=HexColor("#0B1628"), spaceAfter=8)
    h = ParagraphStyle("HK", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=14, leading=17, textColor=HexColor("#007F86"), spaceBefore=8, spaceAfter=4)
    body = ParagraphStyle("BodyK", parent=styles["BodyText"], fontName="Helvetica", fontSize=9.6, leading=13.5, textColor=HexColor("#263746"), spaceAfter=6)
    small = ParagraphStyle("SmallK", parent=body, fontSize=8.3, leading=11.5, textColor=HexColor("#607180"))
    quote = ParagraphStyle("QuoteK", parent=body, fontName="Helvetica-Bold", fontSize=12.5, leading=17, textColor=HexColor("#0B1628"), backColor=HexColor("#DCF9F6"), borderPadding=8, spaceBefore=5, spaceAfter=9)

    doc = SimpleDocTemplate(str(PDF_PATH), pagesize=A4, rightMargin=14*mm, leftMargin=14*mm, topMargin=13*mm, bottomMargin=13*mm)
    story = []
    story.append(Paragraph("KRIYA ZERO", title))
    story.append(Paragraph("The One-Demonstration Physical Skill Compiler", ParagraphStyle("Sub", parent=body, fontName="Helvetica-Bold", fontSize=14, leading=18, textColor=HexColor("#0B1628"))))
    story.append(Paragraph("Smart Education • Team XYRO • iQOO Pune 2026", small))
    story.append(Spacer(1, 5))
    story.append(Paragraph("Show it once. Learn it anywhere. Prove you can do it.", quote))

    story.append(Paragraph("Problem", h))
    story.append(Paragraph("Practical learning is difficult to scale because tutorials can demonstrate a procedure but cannot prove that a learner reproduced it correctly. Labs, ITIs and vocational programs still depend heavily on instructor supervision and manual practical assessment.", body))

    story.append(Paragraph("What we are building", h))
    story.append(Paragraph("KRIYA ZERO is an Android-first, phone-native system that converts a single expert demonstration of a structured physical task into an executable Skill Capsule. The phone observes the workspace, captures learned checkpoints, compiles ordered dependencies, then watches another learner perform the procedure and detects wrong, skipped or out-of-order states. The same representation powers guidance, verification and an evidence-backed skill report.", body))

    story.append(Paragraph("Core innovation", h))
    data = [
        ["Demonstration", "→", "Executable procedural graph", "→", "Live verifier"],
        ["Expert performs once", "", "States + steps + dependencies + checkpoints", "", "PASS / FAIL / sequence error / correction"],
    ]
    table = Table(data, colWidths=[35*mm, 8*mm, 62*mm, 8*mm, 55*mm])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), HexColor("#DCF9F6")),
        ("TEXTCOLOR", (0,0), (-1,-1), HexColor("#0B1628")),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTNAME", (0,1), (-1,1), "Helvetica"),
        ("FONTSIZE", (0,0), (-1,-1), 8.4),
        ("ALIGN", (0,0), (-1,-1), "CENTER"),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("GRID", (0,0), (-1,-1), 0.4, HexColor("#B7DAD8")),
        ("TOPPADDING", (0,0), (-1,-1), 6),
        ("BOTTOMPADDING", (0,0), (-1,-1), 6),
    ]))
    story.append(table)

    story.append(Paragraph("Why it is different", h))
    story.append(Paragraph("Conventional videos teach. Digital SOP and AR work-instruction tools guide procedures that are authored beforehand. KRIYA's wedge is zero-authoring physical verification: a fresh demonstration creates the verifier itself. The hackathon claim is intentionally narrow and defensible — visually observable, structured tabletop procedures rather than arbitrary open-world human activity.", body))

    story.append(Paragraph("Technical approach", h))
    story.append(Paragraph("CameraX provides live phone-camera input. The current prototype learns compact visual checkpoint fingerprints with five-frame stabilization, compiles them into a Skill Capsule, and uses an explicit deterministic Verification Engine for PASS/FAIL and temporal sequence checks. AI/perception supplies evidence; it does not unilaterally decide competence. This makes local embedding or narration models replaceable without rewriting assessment semantics.", body))

    story.append(Paragraph("Current prototype status", h))
    story.append(Paragraph("Implemented: Android + Jetpack Compose shell, live rear-camera flow, Teach → Compile → Verify → Assessment journey, learned visual checkpoints, brightness-normalized similarity, temporal skip detection, correction without restart, evidence-backed metrics, local persistence backend, unit tests, GitHub Actions and CI-generated debug APK. Real-device threshold calibration on the hackathon iQOO phone remains a required validation step.", body))

    story.append(PageBreak())
    story.append(Paragraph("KRIYA ZERO — Phase 1 Execution", title))
    story.append(Paragraph("20-second proof", h))
    story.append(Paragraph("1) Expert demonstrates a 3–5 state tabletop task. 2) KRIYA creates checkpoints. 3) Workspace resets. 4) Learner performs Step 1 correctly. 5) Learner intentionally makes or skips Step 2. 6) KRIYA rejects the state. 7) Learner corrects it without restarting. 8) KRIYA verifies completion and generates the skill report.", body))

    story.append(Paragraph("Who benefits", h))
    impact = [
        ["Engineering labs", "Circuit assembly, instrumentation and practical exercises"],
        ["ITIs / vocational programs", "Scalable repetition and standardized practical assessment"],
        ["Industrial workforce", "Onboarding and procedure practice where expert time is scarce"],
        ["Field technicians", "Guided, measurable procedures without constant remote supervision"],
    ]
    t2 = Table(impact, colWidths=[50*mm, 125*mm])
    t2.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (0,-1), HexColor("#F2F6F9")),
        ("FONTNAME", (0,0), (0,-1), "Helvetica-Bold"),
        ("FONTNAME", (1,0), (1,-1), "Helvetica"),
        ("FONTSIZE", (0,0), (-1,-1), 9),
        ("TEXTCOLOR", (0,0), (-1,-1), HexColor("#263746")),
        ("GRID", (0,0), (-1,-1), 0.4, HexColor("#D8E1E8")),
        ("TOPPADDING", (0,0), (-1,-1), 6),
        ("BOTTOMPADDING", (0,0), (-1,-1), 6),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
    ]))
    story.append(t2)

    story.append(Paragraph("Why iQOO", h))
    story.append(Paragraph("The phone is the product: camera sensing, local perception, offline-critical execution and real-time feedback are central rather than decorative. The hackathon build can use the current deterministic visual layer as the safety net while strengthening perception with a local image-embedding model and local narration on the iQOO device.", body))

    story.append(Paragraph("Feasibility & scope", h))
    story.append(Paragraph("The MVP deliberately focuses on structured tabletop procedures with visually distinguishable state changes. This keeps the 30-hour build testable and demo-safe while preserving a scalable architecture for richer local models, multi-view verification, multilingual guidance and Skill Capsule sharing later.", body))

    story.append(Paragraph("Safe novelty statement", h))
    story.append(Paragraph("KRIYA ZERO learns previously unseen, visually observable, structured physical procedures from a single demonstration and automatically creates reusable guidance and verification checkpoints without manual procedure authoring or task-specific pass/fail training.", quote))

    story.append(Paragraph("Repository", h))
    story.append(Paragraph("https://github.com/omghotekar01-dotcom/KRIYA-ZERO", body))
    story.append(Paragraph("Prototype code, tests, architecture, demo script, on-device acceptance protocol and CI/APK workflow are included in the repository.", small))

    doc.build(story)


if __name__ == "__main__":
    build_deck()
    build_pdf()
    print(PPTX_PATH)
    print(PDF_PATH)
